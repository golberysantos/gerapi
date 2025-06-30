
import tkinter as tk
from tkinter import filedialog, messagebox
import zipfile
import xml.etree.ElementTree as ET
from pptx import Presentation
import os

# Namespaces for WordprocessingML
NS_W = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'

# Função para extrair conteúdo dos controles de conteúdo Rich Text com aliases
def extrair_conteudos_docx(caminho_docx):
    aliases_desejados = {
        'TemaDaApresentacao': 'var_tema_da_apresentacao',
        'ApresentacaoDoProblema': 'var_apresentacao_do_problema',
        'Solucao': 'var_solucao',
        'ModeloDeNegocio': 'var_modelo_de_negocio',
        'PublicoAlvoETamanhoDoMercado': 'var_publico_alvo_e_tamanho_mercado',
        'Encerramento': 'var_encerramento',
        'BenchmarkEConcorrentes': 'var_benchmark_e_concorrentes'
    }
    variaveis = {v: '' for v in aliases_desejados.values()}

    with zipfile.ZipFile(caminho_docx) as docx:
        with docx.open('word/document.xml') as xml_file:
            tree = ET.parse(xml_file)
            root = tree.getroot()
            for sdt in root.iter(NS_W + 'sdt'):
                alias = sdt.find('.//' + NS_W + 'alias')
                if alias is not None:
                    nome_alias = alias.attrib.get(NS_W + 'val')
                    if nome_alias in aliases_desejados:
                        texto = ''.join(node.text or '' for node in sdt.iter(NS_W + 't'))
                        variaveis[aliases_desejados[nome_alias]] = texto.strip()
    return variaveis

# Função para identificar layouts e placeholders
def mapear_layouts(prs):
    layouts = {}
    for layout in prs.slide_layouts:
        nome_layout = layout.name
        placeholders = [ph.name for ph in layout.placeholders]
        layouts[nome_layout] = placeholders
    return layouts

# Função para adicionar slide com layout e preencher placeholders
def adicionar_slide(prs, layout_nome, titulo, conteudo, layouts_map):
    layout = None
    for l in prs.slide_layouts:
        if l.name == layout_nome:
            layout = l
            break
    if layout is None:
        raise ValueError(f"Layout '{layout_nome}' não encontrado.")
    slide = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        if ph.name.startswith("Title") and titulo:
            ph.text = titulo
        elif ph.name.startswith("Content") and conteudo:
            ph.text = conteudo
    return slide

# Função principal para gerar a apresentação
def gerar_apresentacao(caminho_docx, caminho_pptx):
    variaveis = extrair_conteudos_docx(caminho_docx)
    prs = Presentation(caminho_pptx)
    layouts_map = mapear_layouts(prs)

    # Estrutura dos slides
    estrutura = [
        ("Layout_Tema_Padrao", "Title 1", variaveis['var_tema_da_apresentacao']),
        ("Layout_Titulo_e_Conteudo", "Apresentação do Problema", variaveis['var_apresentacao_do_problema']),
        ("Layout_Titulo_e_Conteudo", "Apresentação do Problema", variaveis['var_apresentacao_do_problema']),
        ("Layout_Titulo_e_Conteudo", "Benchmark e Concorrentes", variaveis['var_benchmark_e_concorrentes']),
        ("Layout_Titulo_e_Conteudo", "Público-alvo e Tamanho do mercado", variaveis['var_publico_alvo_e_tamanho_mercado']),
        ("Layout_Titulo_e_Conteudo", "Solução", variaveis['var_solucao']),
        ("Layout_Titulo_e_Conteudo", "Solução", variaveis['var_solucao']),
        ("Layout_Titulo_e_Conteudo", "Modelo de negócio", variaveis['var_modelo_de_negocio']),
        ("Layout_Titulo_e_Conteudo", "Encerramento", variaveis['var_encerramento']),
        ("Layout_Titulo_e_Conteudo", "Redes sociais ou outras Informações", "")
    ]

    for layout_nome, titulo_slide, conteudo in estrutura:
        adicionar_slide(prs, layout_nome, titulo_slide, conteudo, layouts_map)

    prs.save("Apresentacao_ProjetoIntegrador2_Senac.pptx")

# Interface gráfica com Tkinter
def selecionar_docx():
    caminho = filedialog.askopenfilename(filetypes=[("Documentos Word", "*.docx")])
    if caminho:
        entrada_docx.delete(0, tk.END)
        entrada_docx.insert(0, caminho)

def selecionar_pptx():
    caminho = filedialog.askopenfilename(filetypes=[("Apresentações PowerPoint", "*.pptx")])
    if caminho:
        entrada_pptx.delete(0, tk.END)
        entrada_pptx.insert(0, caminho)

def executar():
    caminho_docx = entrada_docx.get()
    caminho_pptx = entrada_pptx.get()
    if not os.path.exists(caminho_docx) or not os.path.exists(caminho_pptx):
        messagebox.showerror("Erro", "Selecione arquivos válidos.")
        return
    try:
        gerar_apresentacao(caminho_docx, caminho_pptx)
        messagebox.showinfo("Sucesso", "Apresentação gerada com sucesso!")
    except Exception as e:
        messagebox.showerror("Erro", str(e))

# Construção da janela
janela = tk.Tk()
janela.title("Gerador de Apresentação - Projeto Integrador 2")

tk.Label(janela, text="Selecione o arquivo Word (.docx):").grid(row=0, column=0, sticky="w")
entrada_docx = tk.Entry(janela, width=60)
entrada_docx.grid(row=1, column=0)
tk.Button(janela, text="Procurar", command=selecionar_docx).grid(row=1, column=1)

tk.Label(janela, text="Selecione o modelo PowerPoint (.pptx):").grid(row=2, column=0, sticky="w")
entrada_pptx = tk.Entry(janela, width=60)
entrada_pptx.grid(row=3, column=0)
tk.Button(janela, text="Procurar", command=selecionar_pptx).grid(row=3, column=1)

tk.Button(janela, text="Gerar Apresentação", command=executar, bg="green", fg="white").grid(row=4, column=0, pady=10)

janela.mainloop()
